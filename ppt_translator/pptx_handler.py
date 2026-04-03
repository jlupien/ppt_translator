"""Read and write PowerPoint files while preserving run-level formatting."""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class RunInfo:
    """A single text run within a paragraph."""
    text: str
    # We don't store formatting — we modify runs in-place on the original slide


@dataclass
class ParagraphInfo:
    """A paragraph containing one or more runs."""
    runs: List[RunInfo] = field(default_factory=list)

    @property
    def text(self) -> str:
        return "".join(r.text for r in self.runs)


@dataclass
class TextFrameInfo:
    """All paragraphs in a text frame, with a reference back to the pptx object."""
    paragraphs: List[ParagraphInfo] = field(default_factory=list)
    text_frame: object = None  # pptx TextFrame reference

    @property
    def full_text(self) -> str:
        return "\n".join(p.text for p in self.paragraphs)


@dataclass
class ShapeInfo:
    """A shape on a slide that contains translatable text."""
    shape_index: int
    text_frames: List[TextFrameInfo] = field(default_factory=list)


@dataclass
class SlideText:
    """All translatable text extracted from a single slide."""
    slide_number: int
    shapes: List[ShapeInfo] = field(default_factory=list)


def _extract_text_frames_from_shape(shape) -> List[TextFrameInfo]:
    """Extract text frame info from a shape, handling tables and groups."""
    frames = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            frames.extend(_extract_text_frames_from_shape(child_shape))
        return frames

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                tf_info = _extract_text_frame(cell.text_frame)
                if tf_info:
                    frames.append(tf_info)
        return frames

    if hasattr(shape, "text_frame"):
        tf_info = _extract_text_frame(shape.text_frame)
        if tf_info:
            frames.append(tf_info)

    return frames


def _extract_text_frame(text_frame) -> Optional[TextFrameInfo]:
    """Extract paragraph and run data from a text frame."""
    paragraphs = []
    has_text = False

    for para in text_frame.paragraphs:
        runs = []
        for run in para.runs:
            runs.append(RunInfo(text=run.text))
            if run.text.strip():
                has_text = True
        if runs:
            paragraphs.append(ParagraphInfo(runs=runs))

    if not has_text:
        return None

    return TextFrameInfo(paragraphs=paragraphs, text_frame=text_frame)


def extract_slide_text(slide, slide_number: int) -> SlideText:
    """Extract all translatable text from a slide."""
    slide_text = SlideText(slide_number=slide_number)

    for shape_index, shape in enumerate(slide.shapes):
        text_frames = _extract_text_frames_from_shape(shape)
        if text_frames:
            slide_text.shapes.append(
                ShapeInfo(shape_index=shape_index, text_frames=text_frames)
            )

    return slide_text


def collect_segments(slide_text: SlideText) -> Tuple[List[str], List[bool]]:
    """Collect all non-empty text segments from a slide for batch translation.

    Multi-run paragraphs are wrapped in <rN> tags so Claude can preserve
    formatting boundaries. Single-run paragraphs are sent as plain text.

    Returns (segments, has_markup) — parallel lists.
    """
    segments = []
    has_markup = []
    for shape in slide_text.shapes:
        for tf in shape.text_frames:
            for para in tf.paragraphs:
                text = para.text
                if not text.strip():
                    continue
                if len(para.runs) > 1:
                    tagged = "".join(
                        f"<r{i+1}>{r.text}</r{i+1}>"
                        for i, r in enumerate(para.runs)
                    )
                    segments.append(tagged)
                    has_markup.append(True)
                else:
                    segments.append(text)
                    has_markup.append(False)
    return segments, has_markup


def apply_translations(
    slide_text: SlideText,
    translated_segments: List[str],
    has_markup: List[bool],
) -> None:
    """Write translated text back into the original pptx text frames.

    Preserves all run-level formatting (font, size, color, bold, italic, etc.)
    by only modifying the .text property of each run.

    For segments with markup tags, parses <rN> tags to map text to the correct
    runs. Falls back to proportional distribution if parsing fails.
    """
    seg_idx = 0

    for shape in slide_text.shapes:
        for tf in shape.text_frames:
            # Build aligned pairs: skip pptx paragraphs that have no runs
            # (empty lines between bullets, etc.) to match our filtered list
            pptx_paras_with_runs = [
                p for p in tf.text_frame.paragraphs if list(p.runs)
            ]

            for para_info, pptx_para in zip(tf.paragraphs, pptx_paras_with_runs):
                original_text = para_info.text
                if not original_text.strip():
                    continue

                if seg_idx >= len(translated_segments):
                    break

                translated = translated_segments[seg_idx]
                markup = has_markup[seg_idx] if seg_idx < len(has_markup) else False
                seg_idx += 1

                pptx_runs = list(pptx_para.runs)

                if markup and len(pptx_runs) > 1:
                    parsed = _parse_tagged_runs(translated, len(pptx_runs))
                    if parsed is not None:
                        for run, text in zip(pptx_runs, parsed):
                            run.text = text
                        continue

                # Fallback: proportional distribution
                _distribute_text_to_runs(pptx_para, para_info.runs, translated)


def _parse_tagged_runs(text: str, num_runs: int) -> Optional[List[str]]:
    """Parse <r1>...</r1><r2>...</r2> tags from translated text.

    Returns a list of run texts if all tags are found in order,
    or None if parsing fails (triggering proportional fallback).
    """
    results = []
    for i in range(1, num_runs + 1):
        pattern = rf"<r{i}>(.*?)</r{i}>"
        match = re.search(pattern, text, re.DOTALL)
        if match is None:
            return None
        results.append(match.group(1))

    if len(results) != num_runs:
        return None

    return results


def _distribute_text_to_runs(pptx_para, run_infos: List[RunInfo], translated: str) -> None:
    """Distribute translated text across the existing runs of a paragraph.

    Strategy:
    - If there's only one run, put all translated text in it.
    - If there are multiple runs, distribute proportionally based on
      original character counts, preserving each run's formatting.
    """
    pptx_runs = list(pptx_para.runs)

    if len(pptx_runs) == 1:
        pptx_runs[0].text = translated
        return

    if len(pptx_runs) == 0:
        return

    # Calculate proportional distribution
    original_lengths = [len(r.text) for r in run_infos]
    total_original = sum(original_lengths)

    if total_original == 0:
        # All runs were empty — put everything in the first run
        pptx_runs[0].text = translated
        for run in pptx_runs[1:]:
            run.text = ""
        return

    # Distribute translated text proportionally across runs
    translated_len = len(translated)
    allocated = 0

    for i, run in enumerate(pptx_runs):
        if i == len(pptx_runs) - 1:
            # Last run gets the remainder
            run.text = translated[allocated:]
        else:
            proportion = original_lengths[i] / total_original
            char_count = round(proportion * translated_len)
            # Try to break at a word boundary
            end = allocated + char_count
            end = _find_word_boundary(translated, end)
            run.text = translated[allocated:end]
            allocated = end


def _find_word_boundary(text: str, pos: int) -> int:
    """Find the nearest word boundary (space) near pos."""
    if pos >= len(text):
        return len(text)
    if pos <= 0:
        return 0

    # Look for a space within 5 characters in either direction
    for offset in range(6):
        if pos + offset < len(text) and text[pos + offset] == " ":
            return pos + offset + 1  # Include the space in the earlier run
        if pos - offset >= 0 and text[pos - offset] == " ":
            return pos - offset + 1

    # No nearby space found — just split at the calculated position
    return pos


def build_deck_summary(prs: Presentation) -> str:
    """Build a compact text summary of the entire presentation for context generation.

    Returns one line per slide with title and body text separated by pipes.
    """
    lines = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            lines.append(f"Slide {slide_number}: {' | '.join(parts)}")
        else:
            lines.append(f"Slide {slide_number}: (no text)")
    return "\n".join(lines)


@dataclass
class ImageInfo:
    """An image shape on a slide."""
    shape: object  # pptx Shape reference
    image_bytes: bytes


def extract_images(slide) -> List[ImageInfo]:
    """Extract all image shapes from a slide."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            if blob:
                images.append(ImageInfo(shape=shape, image_bytes=blob))
    return images


def replace_image(shape, new_image_bytes: bytes) -> None:
    """Replace an image shape's content with new image bytes.

    Preserves the shape's position and size on the slide.
    """
    # Get the image part through the slide's relationship
    pic = shape._element
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    blip = pic.find(f".//{ns}blip")
    rId = blip.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    )
    slide_part = shape.part
    image_part = slide_part.related_part(rId)
    image_part._blob = new_image_bytes


def load_presentation(path: str) -> Presentation:
    """Load a PowerPoint presentation."""
    return Presentation(path)
