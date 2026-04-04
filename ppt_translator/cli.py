"""CLI entry point for the PPT translator."""
from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import List, Sequence

from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

from .pptx_handler import (
    apply_translations,
    build_deck_summary,
    collect_segments,
    extract_images,
    extract_slide_text,
    replace_image,
)
from .translator import DEFAULT_MODEL, TranslationService


def _find_pptx_files(target: Path) -> List[Path]:
    """Find .pptx files at the given path."""
    suffixes = {".pptx"}
    if target.is_file():
        if target.suffix.lower() in suffixes:
            return [target]
        return []
    if target.is_dir():
        return sorted(p for p in target.rglob("*") if p.suffix.lower() in suffixes)
    return []


def _translate_file(
    input_path: Path,
    output_path: Path,
    translator: TranslationService,
    source_lang: str,
    target_lang: str,
    translate_images: bool = False,
) -> None:
    """Translate a single PPTX file."""
    prs = Presentation(str(input_path))
    # Track already-processed images by content hash to skip duplicates
    _processed_image_hashes: set[str] = set()
    total_slides = len(prs.slides)

    # Pass 1: Generate translation context from full deck
    print("  Generating translation guide...", end=" ", flush=True)
    deck_summary = build_deck_summary(prs)
    context = translator.generate_context_summary(
        deck_summary, source_lang, target_lang
    )
    print("done.")

    # Pass 2: Translate each slide with context
    for slide_number, slide in enumerate(prs.slides, start=1):
        print(f"  Slide {slide_number}/{total_slides}...", end=" ", flush=True)

        # Translate text in shapes
        slide_text = extract_slide_text(slide, slide_number)
        segments, has_markup = collect_segments(slide_text)

        text_count = 0
        if segments:
            translated = translator.translate_segments(
                segments, source_lang, target_lang, context=context
            )
            apply_translations(slide_text, translated, has_markup)
            text_count = len(segments)

        # Translate text in images
        img_count = 0
        if translate_images:
            import hashlib
            from .image_handler import translate_image, is_image_too_small

            images = extract_images(slide)
            for img_idx, img_info in enumerate(images):
                # Skip small images (logos, icons, bullets)
                if is_image_too_small(img_info.image_bytes):
                    continue

                # Skip duplicate images (same background reused across slides)
                img_hash = hashlib.sha1(img_info.image_bytes).hexdigest()
                if img_hash in _processed_image_hashes:
                    continue
                _processed_image_hashes.add(img_hash)

                print(f"\n    Processing image {img_idx + 1}/{len(images)}...", end="", flush=True)
                new_bytes = translate_image(
                    img_info.image_bytes, translator,
                    source_lang, target_lang, deck_context=context,
                )
                if new_bytes is not None:
                    replace_image(img_info.shape, new_bytes)
                    img_count += 1
                    print(" done.", end="", flush=True)
                else:
                    print(" skipped.", end="", flush=True)
            if images:
                print(flush=True)  # newline after image progress

        # Progress output
        parts = []
        if text_count:
            parts.append(f"{text_count} segment(s)")
        if img_count:
            parts.append(f"{img_count} image(s)")
        if parts:
            print(f"({', '.join(parts)} translated)")
        else:
            print("(no text)")

    prs.save(str(output_path))
    print(f"  Saved: {output_path}")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="ppt_translator",
        description="Translate PowerPoint slide decks using Claude.",
    )
    parser.add_argument(
        "path",
        help="Path to a .pptx file or directory containing presentations.",
    )
    parser.add_argument(
        "-t",
        "--target",
        required=True,
        help="Target language (e.g., French, Japanese, Spanish).",
    )
    parser.add_argument(
        "-s",
        "--source",
        default="",
        help="Source language. If omitted, Claude will auto-detect.",
    )
    parser.add_argument(
        "-o",
        "--output",
        default=None,
        help="Output file path. Only valid with a single input file.",
    )
    parser.add_argument(
        "--model",
        default=DEFAULT_MODEL,
        help=f"Claude model to use (default: {DEFAULT_MODEL}).",
    )
    parser.add_argument(
        "--translate-images",
        action="store_true",
        help="Also translate text within images (requires easyocr).",
    )
    parser.add_argument(
        "--debug-images",
        action="store_true",
        help="Write detailed image translation logs to ~/.ppt_translator/logs/.",
    )
    return parser


def main(argv: Sequence[str] | None = None) -> None:
    parser = build_parser()
    args = parser.parse_args(argv)

    target_path = Path(args.path).expanduser().resolve()
    files = _find_pptx_files(target_path)

    if not files:
        print(f"No .pptx files found at: {target_path}")
        sys.exit(1)

    if args.output and len(files) > 1:
        print("Error: --output can only be used with a single input file.")
        sys.exit(1)

    if args.debug_images:
        from .image_handler import enable_debug_logging
        log_path = enable_debug_logging()
        print(f"Image debug log: {log_path}")

    try:
        translator = TranslationService(model=args.model)
    except RuntimeError as e:
        print(f"Error: {e}")
        sys.exit(1)

    # If input is a directory, output goes to a "translated/" subdirectory
    is_directory_mode = target_path.is_dir()

    for pptx_file in files:
        if args.output:
            out = Path(args.output).expanduser().resolve()
        elif is_directory_mode:
            out_dir = target_path / "translated"
            out_dir.mkdir(exist_ok=True)
            out = out_dir / pptx_file.name
        else:
            out = pptx_file.with_name(f"{pptx_file.stem}_translated{pptx_file.suffix}")

        print(f"Translating: {pptx_file.name}")
        source_label = args.source or "auto-detect"
        print(f"  {source_label} → {args.target}")

        _translate_file(
            pptx_file, out, translator, args.source, args.target,
            translate_images=args.translate_images,
        )

    print("Done.")
