"""Create a test PPTX with various formatting for translation testing."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1: Title slide with plain text
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Quarterly Business Report"
slide1.placeholders[1].text = "Prepared by the Analytics Team"

# Slide 2: Mixed formatting (bold + italic + colored runs)
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Key Findings"

body = slide2.placeholders[1].text_frame
body.clear()

p = body.paragraphs[0]
run1 = p.add_run()
run1.text = "Revenue grew by "
run1.font.size = Pt(18)

run2 = p.add_run()
run2.text = "25 percent"
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.color.rgb = RGBColor(0x00, 0x80, 0x00)  # Green

run3 = p.add_run()
run3.text = " compared to last quarter."
run3.font.size = Pt(18)

p2 = body.add_paragraph()
run4 = p2.add_run()
run4.text = "Customer satisfaction "
run4.font.size = Pt(18)

run5 = p2.add_run()
run5.text = "improved significantly"
run5.font.size = Pt(18)
run5.font.italic = True
run5.font.color.rgb = RGBColor(0x00, 0x00, 0xCC)  # Blue

run6 = p2.add_run()
run6.text = " across all regions."
run6.font.size = Pt(18)

# Slide 3: Table
slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
slide3.shapes.add_table(
    rows=3, cols=2, left=Inches(1), top=Inches(1.5), width=Inches(8), height=Inches(3)
).table

table = slide3.shapes[-1].table
cells = [
    ("Region", "Revenue"),
    ("North America", "Twelve million dollars"),
    ("Europe", "Eight million dollars"),
]
for row_idx, (c1, c2) in enumerate(cells):
    table.cell(row_idx, 0).text = c1
    table.cell(row_idx, 1).text = c2
    if row_idx == 0:
        for col in range(2):
            for p in table.cell(row_idx, col).text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True

# Slide 4: Multiple text boxes
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
txBox1 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
txBox1.text_frame.text = "Our mission is to deliver excellence."

txBox2 = slide4.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
p = txBox2.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Innovation drives our success."
run.font.size = Pt(24)
run.font.color.rgb = RGBColor(0x80, 0x00, 0x80)  # Purple

prs.save("test_presentation.pptx")
print("Created test_presentation.pptx with 4 slides")
